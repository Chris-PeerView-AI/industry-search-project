import os
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

# Paths
PROJECT_ROOT = os.path.dirname(os.path.dirname(__file__))
MODULES_DIR = os.path.join(PROJECT_ROOT, "modules")
template_path = os.path.join(MODULES_DIR, "downloaded_summary_template.pptx")
map_image_path = os.path.join(PROJECT_ROOT, "sample.png")
output_path = os.path.join(MODULES_DIR, "test_output_summary_slide.pptx")

# Load template
prs = Presentation(template_path)
slide = prs.slides[0]

# Placeholder replacements
replacements = {
    "{TBD TITLE}": "Exhibit B: Industry Summary – Austin Doggie Daycare",
    "{TBD AS OF DATE}": "July 31, 2025",
    "{TBD TOTAL BUSINESSES}": "32",
    "{TBD TRUSTED BUSINESSES}": "24",
    "{TBD: MEAN REVENUE}": "$418,000",
    "{TBD YOY GROWTH}": "+6.3%",
    "{TBD MEDIAM REVENUE}": "$390,500",
    "{TBD AVERAGE TICKET SIZE}": "$23.10",
    "{TBD SUMMARY ANALYSIS}": (
        "The Austin pet care market continues to show strong growth, led by rising "
        "customer demand and steady repeat visitation patterns."
    )
}

# Replace placeholders in all shapes
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    full_text = shape.text
    updated_text = full_text
    for placeholder, value in replacements.items():
        updated_text = updated_text.replace(placeholder, value)

    if updated_text != full_text:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = updated_text

        # Font logic
        if "{TBD TITLE}" in full_text:
            run.font.name = "Montserrat"
            run.font.size = Pt(36)
        else:
            run.font.name = "Arial"
            run.font.size = Pt(11)

        print(f"✅ Replaced text in shape:\n{full_text[:30]}...")

# Insert map image: shift left by 1 inch, down by 1 inch, keep right edge the same
with Image.open(map_image_path) as img:
    img_width_px, img_height_px = img.size
    dpi = 96
    original_width_in = 2.5  # assumed original width
    final_width_in = original_width_in + 1.0  # extend left by 1 inch
    scale = (final_width_in * dpi) / img_width_px
    final_height_in = (img_height_px * scale) / dpi

# Slide dimensions
slide_width_in = prs.slide_width.inches
right_edge_in = slide_width_in - 0.5  # original 0.5" right margin
left_edge_in = right_edge_in - final_width_in
top_edge_in = 1.0 + 1.0  # moved down 1 inch

# Add image
slide.shapes.add_picture(
    map_image_path,
    Inches(left_edge_in), Inches(top_edge_in),
    width=Inches(final_width_in),
    height=Inches(final_height_in)
)
print(f"✅ Inserted image from x={left_edge_in:.2f}in to x={right_edge_in:.2f}in")

# Save file
prs.save(output_path)
print(f"✅ Final saved to: {output_path}")
