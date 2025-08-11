# --- File: modules/map_generator.py (PATCH v3.1) ---
# Fixes JS injection braces to avoid SyntaxError in f-strings and ensures proper map centering/aspect ratio.

from __future__ import annotations

import math
import os
import time
from dataclasses import dataclass
from typing import Iterable, Tuple, Optional
from string import Template

import pandas as pd
import folium
from branca.element import Element
from geopy.distance import geodesic

from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.common.exceptions import JavascriptException
from selenium.webdriver.common.by import By

POSI_TILE_URL = "https://{s}.basemaps.cartocdn.com/light_all/{z}/{x}/{y}{r}.png"
POSI_TILE_ATTR = "© OpenStreetMap contributors © CARTO"
WINDOW_DEFAULT = (1200, 800)
TILE_WAIT_HARD_TIMEOUT_SEC = 12.0
TILE_POLL_INTERVAL_SEC = 0.25


@dataclass
class MapMeta:
    center_lat: float
    center_lng: float
    radius_m: float
    count: int
    desired_zoom: float


def compute_zoom_for_circle(lat_deg: float, radius_m: float, window_h_px: int, *, target_fraction: float) -> float:
    rpx_target = (window_h_px * target_fraction) / 2.0
    numerator = 156543.03392 * math.cos(math.radians(lat_deg)) * rpx_target
    if radius_m <= 0:
        radius_m = 200.0
    z = math.log2(max(numerator / radius_m, 1e-6))
    return max(1.0, min(z, 19.0))


def _bbox_midpoint(df: pd.DataFrame) -> tuple[float, float]:
    return (float((df["latitude"].min() + df["latitude"].max()) / 2.0),
            float((df["longitude"].min() + df["longitude"].max()) / 2.0))


def _radius_from_center(center: tuple[float, float], df: pd.DataFrame) -> int:
    clat, clng = center
    return max(200, int(max(
        geodesic((clat, clng), (lat, lng)).km for lat, lng in zip(df["latitude"], df["longitude"])) * 1000))


def build_map(df: pd.DataFrame, *, zoom_fraction: float = 0.75, window: Tuple[int, int] = WINDOW_DEFAULT) -> tuple[
    folium.Map, MapMeta]:
    df = df[df["latitude"].notna() & df["longitude"].notna()].copy()
    if df.empty:
        raise ValueError("No valid lat/lon rows")
    center_lat, center_lng = _bbox_midpoint(df)
    radius_m = _radius_from_center((center_lat, center_lng), df)
    desired_zoom = compute_zoom_for_circle(center_lat, radius_m, window[1], target_fraction=zoom_fraction)

    m = folium.Map(location=[center_lat, center_lng], zoom_start=int(round(desired_zoom)),
                   tiles=None, control_scale=False, zoom_control=False, max_zoom=19,
                   width=window[0], height=window[1])
    folium.TileLayer(POSI_TILE_URL, name="Positron", attr=POSI_TILE_ATTR, control=False).add_to(m)

    m.get_root().html.add_child(Element(
        """<style>.leaflet-container .leaflet-tile { opacity: 0.92; filter: saturate(0.85) brightness(1.02); }</style>"""
    ))

    css_tpl = Template("""
    <style>
      html, body { margin:0; padding:0; overflow:hidden; }
      #$map_id { position:fixed; inset:0; width:100vw; height:100vh; }
      .folium-map, .leaflet-container { width:100vw !important; height:100vh !important; }
    </style>
    """)
    m.get_root().html.add_child(Element(css_tpl.substitute(map_id=m.get_name())))

    m.get_root().html.add_child(Element(
        """<style>.leaflet-container .leaflet-tile { opacity: 0.92; filter: saturate(0.85) brightness(1.02); }</style>"""))
    js_tpl = Template("""
       <script>
       (function(){
         var m=$map_id;
         if(m){
           m.options.zoomSnap = 0; m.options.zoomDelta = 0.1;
           m.setZoom($zoom, {animate:false});
           m.setView([$lat, $lng], $zoom);
           setTimeout(function(){ m.invalidateSize(false); }, 60);
         }
       })();
       </script>
       """)
    js_str = js_tpl.substitute(
        map_id=m.get_name(),
        zoom=f"{desired_zoom:.4f}",
        lat=f"{center_lat:.6f}",
        lng=f"{center_lng:.6f}",
    )
    m.get_root().html.add_child(Element(js_str))
    folium.Circle(location=[center_lat, center_lng], radius=radius_m, color="#2c7fb8", fill=True,
                  fill_opacity=0.05, weight=1.2, opacity=0.8, dash_array="6 6").add_to(m)
    for _, row in df.iterrows():
        fill = "#2ca25f" if str(row.get("benchmark", "")).lower() == "trusted" else "#7f8c8d"
        folium.CircleMarker(location=[row["latitude"], row["longitude"]], radius=8, weight=2,
                            color="#ffffff", fill=True, fill_color=fill, fill_opacity=0.95).add_to(m)

    m.get_root().html.add_child(Element(
        """<div style='position: fixed; top: 12px; right: 12px; z-index: 9999; background: rgba(255,255,255,0.96); border: 1px solid #d0d0d0; border-radius: 6px; padding: 14px 16px; font-size: 18px; line-height: 1.6; box-shadow: 0 2px 8px rgba(0,0,0,.15);'><div style='font-weight:600; margin-bottom:8px;'>Legend</div><div><span style='display:inline-block;width:14px;height:14px;border:2px solid #fff;background:#2ca25f;border-radius:50%;margin-right:8px;'></span>Trusted</div><div><span style='display:inline-block;width:14px;height:14px;border:2px solid #fff;background:#7f8c8d;border-radius:50%;margin-right:8px;'></span>Other</div></div>"""))
    return m, MapMeta(center_lat, center_lng, float(radius_m), int(len(df)), float(desired_zoom))


def save_html_and_png(m: folium.Map, html_path: str, png_path: str, window: Tuple[int, int] = WINDOW_DEFAULT) -> None:
    m.save(html_path)
    options = Options()
    options.add_argument("--headless=new")
    options.add_argument("--no-sandbox")
    options.add_argument("--disable-dev-shm-usage")
    options.add_argument("--force-device-scale-factor=1")
    driver = webdriver.Chrome(options=options)
    try:
        driver.set_window_size(*window)
        driver.get("file://" + os.path.abspath(html_path))
        deadline = time.time() + TILE_WAIT_HARD_TIMEOUT_SEC
        driver.execute_script("document.body.style.overflow='hidden'")
        while time.time() < deadline:
            try:
                loading = driver.execute_script("return document.querySelectorAll('.leaflet-tile-loading').length")
                tiles_seen = driver.execute_script(
                    "return document.querySelectorAll('.leaflet-tile, .leaflet-tile-loaded').length")
            except JavascriptException:
                loading, tiles_seen = 0, 0
            if int(loading) == 0 and int(tiles_seen) > 0:
                time.sleep(0.4)
                break
            time.sleep(TILE_POLL_INTERVAL_SEC)
        elem = driver.find_element(By.ID, m.get_name())
        elem.screenshot(png_path)
    finally:
        driver.quit()


def generate_map_png_from_summaries(summaries: Iterable[dict], output_path: str, *, zoom_fraction: float = 0.75,
                                    aspect_ratio: Optional[float] = None, window_height_px: int = 800) -> bool:
    df = pd.DataFrame(list(summaries))
    df = df[df["latitude"].notna() & df["longitude"].notna()]
    if df.empty:
        return False
    if aspect_ratio is None:
        aspect_ratio = 3 / 2
    window = (int(window_height_px * aspect_ratio), int(window_height_px))
    m, _ = build_map(df, zoom_fraction=zoom_fraction, window=window)
    html_path = output_path.replace(".png", ".html")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    save_html_and_png(m, html_path, output_path, window=window)
    return True


def generate_map_png_from_project(project_id: str, supabase, output_dir: str, *, zoom_fraction: float = 0.75,
                                  aspect_ratio: Optional[float] = None) -> Optional[str]:
    resp = supabase.table("enigma_summaries").select("name, latitude, longitude, benchmark").eq("project_id",
                                                                                                project_id).execute()
    df = pd.DataFrame(resp.data or [])
    df = df[df["latitude"].notna() & df["longitude"].notna()]
    if df.empty:
        return None
    if aspect_ratio is None:
        aspect_ratio = 3 / 2
    window = (int(800 * aspect_ratio), 800)
    m, _ = build_map(df, zoom_fraction=zoom_fraction, window=window)
    html_path = os.path.join(output_dir, "test_map.html")
    png_path = os.path.join(output_dir, "test_map.png")
    save_html_and_png(m, html_path, png_path, window=window)
    return png_path


def _find_named(slide, *names):
    for shp in slide.shapes:
        if getattr(shp, "name", "") in names:
            return shp
    return None


def _chart_anchor_dims_from_template(template_path: str):
    from pptx import Presentation
    from pptx.util import Inches
    ppt = Presentation(template_path)
    slide = ppt.slides[0]
    anchor = _find_named(slide, "ChartAnchor", "Chart", "ImageAnchor")
    if anchor:
        return anchor.left, anchor.top, anchor.width, anchor.height
    max_area, best = 0, None
    for shp in slide.shapes:
        try:
            area = shp.width * shp.height
            if area > max_area:
                max_area, best = area, shp
        except Exception:
            pass
    if best:
        return best.left, best.top, best.width, best.height
    left = Inches(0.75)
    top = Inches(1.2)
    width = ppt.slide_width - 2 * left
    height = Inches(4.0)
    return left, top, width, height


def generate_map_chart(output_path, summaries):
    try:
        from modules.slides_exhibit import EXHIBIT_TEMPLATE
    except Exception:
        EXHIBIT_TEMPLATE = "modules/downloaded_exhibit_template.pptx"
    left, top, width, height = _chart_anchor_dims_from_template(EXHIBIT_TEMPLATE)
    aspect_ratio = float(width) / float(height) if height else (3 / 2)
    return generate_map_png_from_summaries(summaries, output_path, zoom_fraction=0.75, aspect_ratio=aspect_ratio,
                                           window_height_px=800)
