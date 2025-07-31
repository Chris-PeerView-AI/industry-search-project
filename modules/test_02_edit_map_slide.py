import os
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

# Paths
PROJECT_ROOT = os.path.dirname(os.path.dirname(__file__))
map_image_path = os.path.join(PROJECT_ROOT, "sample.png")

# Load presentation and slide
prs = Presentation("downloaded_template.pptx")
slide = prs.slides[0]

# Replace full shape text content if it matches target
def replace_text_full(shape, target_text, new_text, font_name, font_size_pt):
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    text_found = False

    if target_text in tf.text:
        tf.clear()  # remove all paragraphs
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = new_text
        run.font.name = font_name
        run.font.size = Pt(font_size_pt)
        return True

    return False


# Replace title and analysis text
title_done = False
analysis_done = False
for shape in slide.shapes:
    if not title_done:
        title_done = replace_text_full(
            shape,
            "Exhibit {TBD}: Map Overview",
            "Exhibit A: Map Overview",
            font_name="Montserrat",
            font_size_pt=36
        )
        if title_done:
            print("✅ Title replaced with Montserrat 36pt (fallback Roboto)")

    if not analysis_done:
        analysis_done = replace_text_full(
            shape,
            "{To be Replace}",
            "This map shows the location of all trusted and untrusted businesses in the Austin doggie daycare market.",
            font_name="Arial",
            font_size_pt=11
        )
        if analysis_done:
            print("✅ Analysis text replaced with Arial 11pt")


# Insert map image, centered and scaled
with Image.open(map_image_path) as img:
    img_width_px, img_height_px = img.size
    target_width_in = 7.5
    dpi = 96
    scale = (target_width_in * dpi) / img_width_px
    target_height_in = (img_height_px * scale) / dpi

slide_width_in = prs.slide_width.inches
img_x = (slide_width_in - target_width_in) / 2

slide.shapes.add_picture(
    map_image_path,
    Inches(img_x), Inches(2.0),
    width=Inches(target_width_in),
    height=Inches(target_height_in)
)
print(f"✅ Map inserted at x={img_x:.2f}in, width=7.5in")

# Save output
prs.save("test_output_slide.pptx")
print("✅ Saved final version as test_output_slide.pptx")
