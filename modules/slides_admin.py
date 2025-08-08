# slides_admin.py — merged version (fixed overlap)
# Formatting-preserving replacements + vertical centering + font scaling + optional auto-generated cover art
# Now detects a separate DATE textbox and positions it **below** the Industry/City block to avoid overlap.

from datetime import datetime
import os
from typing import Dict, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

try:
    from PIL import Image, ImageDraw
    _PIL_AVAILABLE = True
except ImportError:
    _PIL_AVAILABLE = False

BRAND_BLUE = (45, 125, 210)
BRAND_GREEN = (0, 150, 136)

# ---------------------------------
# Formatting-preserving replacement
# ---------------------------------

def _replace_placeholders_in_shape(shape, replacements: Dict[str, str]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    replaced_any = False
    for p in tf.paragraphs:
        for run in p.runs:
            if not run.text:
                continue
            original = run.text
            new_text = original
            for k, v in replacements.items():
                if k in new_text:
                    new_text = new_text.replace(k, v)
            if new_text != original:
                run.text = new_text
                replaced_any = True
    full_text = "".join(p.text for p in tf.paragraphs)
    if any(k in full_text for k in replacements.keys()) and not replaced_any:
        for k, v in replacements.items():
            full_text = full_text.replace(k, v)
        if tf.paragraphs and tf.paragraphs[0].runs:
            for p in tf.paragraphs:
                for r in p.runs:
                    r.text = ""
            tf.paragraphs[0].runs[0].text = full_text

# -----------------------
# Layout helper functions
# -----------------------

def _scale_font(base_size: int, text: str, soft_limit: int, hard_min: int) -> int:
    if not text:
        return base_size
    if len(text) <= soft_limit:
        return base_size
    over = len(text) - soft_limit
    return max(hard_min, base_size - (over // 4))


def _compute_title_block_top(slide_height, lines_count: int, avg_line_pt: int):
    top_margin = Inches(1.7)   # below banner
    bottom_margin = Inches(2.0)
    content_height = Pt(avg_line_pt * max(1, lines_count) * 1.1)
    available = slide_height - (top_margin + bottom_margin)
    if content_height >= available:
        return int(top_margin)
    offset = int((available - content_height) / 2)
    return int(top_margin + offset)


def _make_cover_image(output_dir: str, industry: str, width_px: int = 1700, height_px: int = 2200):
    if not _PIL_AVAILABLE:
        return None
    os.makedirs(output_dir, exist_ok=True)
    path = os.path.join(output_dir, "cover_art.png")
    img = Image.new("RGB", (width_px, height_px), BRAND_BLUE)
    draw = ImageDraw.Draw(img)
    for y in range(height_px):
        t = y / max(1, height_px - 1)
        r = int(BRAND_BLUE[0] * (1 - t) + BRAND_GREEN[0] * t)
        g = int(BRAND_BLUE[1] * (1 - t) + BRAND_GREEN[1] * t)
        b = int(BRAND_BLUE[2] * (1 - t) + BRAND_GREEN[2] * t)
        draw.line([(0, y), (width_px, y)], fill=(r, g, b))
    for i in range(-height_px, width_px, 120):
        draw.line([(i, 0), (i + height_px, height_px)], fill=(255, 255, 255), width=2)
    overlay = Image.new("RGBA", (width_px, height_px), (255, 255, 255, 210))
    img = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")
    wm = (industry or "").upper()[:48]
    if wm:
        # center watermark roughly lower-middle; Pillow default font used for portability
        text_w = draw.textlength(wm)
        x = (width_px - text_w) // 2
        y = int(height_px * 0.62)
        draw.text((x, y), wm, fill=(255, 255, 255))
    img.save(path, "PNG", optimize=True)
    return path

# --------------
# Public API
# --------------

def generate_title_slide(
    project_output_dir: str,
    template_path: str = "modules/downloaded_title_template.pptx",
    city: str = "City, ST",
    industry: str = "Industry",
    date_str: Optional[str] = None,
    subtitle: Optional[str] = None,
    add_cover_art: bool = True,
) -> str:
    if date_str is None:
        date_str = datetime.now().strftime("%B %Y")

    os.makedirs(project_output_dir, exist_ok=True)
    out_path = os.path.join(project_output_dir, "slide_1_title.pptx")

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Optional background art
    if add_cover_art:
        art_path = _make_cover_image(project_output_dir, industry)
        if art_path:
            pic = slide.shapes.add_picture(art_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

    # Replace placeholders
    replacements = {
        "{TBD INDUSTRY}": industry,
        "{TBD LOCATION}": city,
        "{TBD DATE}": date_str,
    }
    if subtitle is not None:
        replacements["{TBD SUBTITLE}"] = subtitle

    for shape in slide.shapes:
        _replace_placeholders_in_shape(shape, replacements)

    # Identify the separate title and date shapes
    title_shape = None
    date_shape = None
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text_all = "\n".join([run.text for p in shape.text_frame.paragraphs for run in p.runs])
        if (industry in text_all) and (city in text_all):
            title_shape = shape
        if date_str in text_all:
            date_shape = shape

    # Scale & center title block
    if title_shape is not None:
        tf = title_shape.text_frame
        non_empty = [p for p in tf.paragraphs if any(r.text.strip() for r in p.runs)]
        if len(non_empty) >= 1:
            for r in non_empty[0].runs:  # Industry line
                base = int((r.font.size or Pt(54)).pt)
                r.font.size = Pt(_scale_font(base, industry, soft_limit=24, hard_min=28))
                r.font.bold = True if r.font.bold is None else r.font.bold
        if len(non_empty) >= 2:
            for r in non_empty[1].runs:  # City line
                base = int((r.font.size or Pt(40)).pt)
                r.font.size = Pt(_scale_font(base, city, soft_limit=28, hard_min=24))
        # paragraph alignment safeguard
        for p in tf.paragraphs:
            p.alignment = p.alignment or PP_ALIGN.CENTER
        # vertical centering calc
        sizes = []
        for p in non_empty[:3]:
            for r in p.runs:
                if r.font.size:
                    sizes.append(r.font.size.pt)
        avg_pt = int(sum(sizes) / len(sizes)) if sizes else 36
        title_shape.top = _compute_title_block_top(prs.slide_height, lines_count=min(3, len(non_empty)), avg_line_pt=avg_pt)
        title_shape.left = int((prs.slide_width - title_shape.width) / 2)

    # Place date shape **below** the title to avoid overlap
    if date_shape is not None and title_shape is not None:
        GAP = Inches(0.35)
        date_shape.left = int((prs.slide_width - date_shape.width) / 2)  # keep centered
        date_shape.top = int(title_shape.top + title_shape.height + GAP)
        # Ensure date text is centered too
        for p in date_shape.text_frame.paragraphs:
            p.alignment = p.alignment or PP_ALIGN.CENTER

    prs.save(out_path)
    print(f"✅ Saved title slide to: {out_path}")
    return out_path

__all__ = ["generate_title_slide", "_replace_placeholders_in_shape"]
