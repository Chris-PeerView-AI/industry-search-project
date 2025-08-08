# slides_summary.py

import os
import subprocess
from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw

SUMMARY_TEMPLATE = "modules/downloaded_summary_template.pptx"
INDIVIDUAL_TEMPLATE = "modules/downloaded_businessview_template.pptx"


def _replace_text_preserve_format(shape, replacements: dict):
    """Replace placeholder substrings within runs so template formatting stays intact."""
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            if not r.text:
                continue
            txt = r.text
            for k, v in replacements.items():
                if k in txt:
                    txt = txt.replace(k, v)
            r.text = txt


def generate_summary_slide(output_path, trusted, end_date, summary_stats, summary_analysis,
                           city: str = "", industry: str = "", map_image_path: str | None = None):
    ppt = Presentation(SUMMARY_TEMPLATE)
    slide = ppt.slides[0]

    # 1) Text replacements only (no hard-coded fonts for titles/stats)
    replacements = {
        "{TBD TITLE}": f"{city}: {industry}".strip(": "),
        "{TBD AS OF DATE}": end_date,
        "{TBD TOTAL BUSINESSES}": str(summary_stats.get("total", "-")),
        "{TBD TRUSTED BUSINESSES}": str(summary_stats.get("trusted", "-")),
        "{TBD: MEAN REVENUE}": f"${summary_stats.get('mean_revenue', 0):,.0f}",
        "{TBD YOY GROWTH}": f"{summary_stats.get('mean_yoy', 0):.1f}%",
        "{TBD MEDIAN REVENUE}": f"${summary_stats.get('median_revenue', 0):,.0f}",
        "{TBD MEDIAM REVENUE}": f"${summary_stats.get('median_revenue', 0):,.0f}",  # legacy typo support
        "{TBD AVERAGE TICKET SIZE}": f"${summary_stats.get('avg_ticket', 0):,.0f}",
        "{TBD MEDIAN TICKET}": f"${summary_stats.get('median_ticket', 0):,.0f}",
    }

    for shape in slide.shapes:
        _replace_text_preserve_format(shape, replacements)

    # 1b) Normalize spacing in the stats/overview box if present (tighten leading and paragraph spacing)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            txt = shape.text_frame.text or ""
            if any(k in txt for k in ["Total", "Trusted", "$", "%"]):
                for para in shape.text_frame.paragraphs:
                    para.line_spacing = 1.1
                    para.space_after = Pt(2)

    # 2) Summary body block (replace a single placeholder) with smaller font + tighter spacing
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and "{TBD SUMMARY ANALYSIS}" in (shape.text or ""):
            tf = shape.text_frame
            tf.clear()

            # Split analysis into paragraphs for consistent spacing
            paragraphs = [p.strip() for p in summary_analysis.strip().split("\n\n") if p.strip()]
            if not paragraphs:
                paragraphs = [summary_analysis.strip()]

            # First paragraph
            p = tf.paragraphs[0]
            p.text = paragraphs[0]
            p.line_spacing = 1.1
            p.space_after = Pt(2)
            for run in p.runs:
                run.font.size = Pt(8)

            # Remaining paragraphs
            for chunk in paragraphs[1:]:
                pn = tf.add_paragraph()
                pn.text = chunk
                pn.line_spacing = 1.1
                pn.space_after = Pt(2)
                for run in pn.runs:
                    run.font.size = Pt(8)
            break

    # 3) Map image: prefer a named anchor; else fallback to right panel with safe scaling
    anchor = None
    for shp in slide.shapes:
        if getattr(shp, "name", "").lower() == "mapplaceholder" or (getattr(shp, "has_text_frame", False) and "{TBD MAP}" in (shp.text or "")):
            anchor = shp
            break

    if map_image_path and os.path.exists(map_image_path):
        if anchor:
            slide.shapes.add_picture(map_image_path, anchor.left, anchor.top, width=anchor.width, height=anchor.height)
            if getattr(anchor, "has_text_frame", False):
                anchor.text_frame.clear()
        else:
            img = Image.open(map_image_path)
            img_w, img_h = img.size  # pixels
            dpi = 96.0
            img_w_in, img_h_in = img_w / dpi, img_h / dpi
            max_w_in, max_h_in = 4.2, 3.7
            ratio = min(max_w_in / img_w_in, max_h_in / img_h_in)
            final_w = Inches(img_w_in * ratio)
            final_h = Inches(img_h_in * ratio)
            slide.shapes.add_picture(
                map_image_path,
                left=Inches(4.0),
                top=Inches(2.25),
                width=final_w,
                height=final_h,
            )

    ppt.save(output_path)
    print(f"✅ Saved summary slide to: {output_path}")


def generate_llama_summary(slide_summaries: dict, model_name: str = "llama3") -> str:
    sentiment = "neutral"
    try:
        mean_yoy = float(slide_summaries.get("yoy", "0").split("Avg:")[-1].split("%")[0])
        if mean_yoy > 10:
            sentiment = "positive"
        elif mean_yoy < 0:
            sentiment = "negative"
    except Exception:
        pass

    prompt = f"""
You are a market research consultant. Based on the following summaries:

1. Revenue: {slide_summaries.get('revenue')}
2. YoY Growth: {slide_summaries.get('yoy')}
3. Ticket Size: {slide_summaries.get('ticket')}
4. Market Size: {slide_summaries.get('market')}

Write a professional, approximately 700-word summary about the market's attractiveness, notable patterns, standout businesses, and whether this is a good location to open a new business.
The tone should be {sentiment}. If growth is above 10%, highlight strong momentum. If it's below 0%, note concerning trends. If in-between, maintain a balanced tone.
""".strip()

    result = subprocess.run(
        ["ollama", "run", model_name],
        input=prompt.encode("utf-8"),
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
    )
    return result.stdout.decode("utf-8").strip()


def get_market_size_analysis():
    return (
        "This chart compares the total verified revenue from businesses with high-quality data to an estimate "
        "for the full local market. The upper bound assumes that businesses without usable data perform similarly "
        "to those with data, which likely overstates true market size. Poor data quality is often associated with smaller "
        "businesses or those facing operational challenges."
    )


def generate_individual_business_slide(output_path, business: dict, end_date: str, industry: str, city: str):
    print(f"🧩 Generating business slide for: {business.get('name')}")

    ppt = Presentation(INDIVIDUAL_TEMPLATE)
    slide = ppt.slides[0]

    replacements = {
        "{TBD TITLE}": business.get("name", "Business"),
        "{TBD AS OF DATE}": end_date,
        "{TBD ADDRESS}": f"{business.get('address', '')}",
        "{TBD: MEAN REVENUE}": f"${business.get('annual_revenue', 0):,.0f}",
        "{TBD YOY GROWTH}": f"{business.get('yoy_growth', 0)*100:.1f}%",
        "{TBD AVERAGE TICKET SIZE}": f"${business.get('ticket_size', 0):,.0f}",
    }

    summary_text = (
        f"{business.get('name')} reported approximately ${business.get('annual_revenue', 0):,.0f} in annual revenue "
        f"as of {end_date}. The average ticket size was ${business.get('ticket_size', 0):,.0f}, and year-over-year revenue "
        f"{'increased' if business.get('yoy_growth', 0) >= 0 else 'declined'} by {abs(business.get('yoy_growth', 0) * 100):.1f}%."
    )
    if business.get("tier_reason"):
        summary_text += f" This business was included due to its: {business['tier_reason'].strip()}."

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        text = shape.text_frame.text
        handled = False

        for key, val in replacements.items():
            if key == "{TBD TITLE}" and key in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = val
                run.font.name = "Montserrat"
                run.font.size = Pt(30)
                run.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                handled = True
                break
            elif key in text:
                text = text.replace(key, val)

        if handled:
            continue

        if "{TBD SUMMARY ANALYSIS}" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = summary_text
            run.font.size = Pt(7)
            continue

        shape.text_frame.text = text
        for p in shape.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT

    # Optional: map pin overlay if lat/lon is available
    if business.get("latitude") and business.get("longitude"):
        try:
            map_path = f"modules/output_map_with_pin_{business['id']}.png"
            base_map = Image.open(f"modules/output/{business['project_id']}/slide_25_map.png").convert("RGBA")
            draw = ImageDraw.Draw(base_map)
            x, y = business['map_x'], business['map_y']  # assuming preprojected pixel coords
            draw.ellipse((x - 6, y - 6, x + 6, y + 6), fill="gold", outline="black")
            base_map.save(map_path)

            slide.shapes.add_picture(map_path, Inches(4.0), Inches(2.25), width=Inches(4.0), height=Inches(3.5))
        except Exception as e:
            print(f"⚠️ Map rendering error: {e}")

    ppt.save(output_path)
    print(f"✅ Saved individual business slide: {output_path}")


def get_latest_period_end(supabase: object, project_id: str) -> str:
    resp = (
        supabase.table("enigma_metrics")
        .select("period_end_date")
        .eq("project_id", project_id)
        .order("period_end_date", desc=True)
        .limit(1)
        .execute()
    )
    if resp.data:
        return datetime.strptime(resp.data[0]["period_end_date"], "%Y-%m-%d").strftime("%B %Y")
    return datetime.now().strftime("%B %Y")


# -----------------------
# Appendix: business table
# -----------------------

def _find_anchor(slide, name_lower: str):
    for shp in slide.shapes:
        if getattr(shp, "name", "").lower() == name_lower:
            return shp
    return None


def generate_paginated_business_table_slides(output_dir: str, businesses: list, base_title: str):
    from math import ceil

    rows_per_slide = 15
    total_slides = max(1, ceil(len(businesses) / rows_per_slide))

    for i in range(total_slides):
        batch = businesses[i * rows_per_slide:(i + 1) * rows_per_slide]
        ppt = Presentation("modules/downloaded_businesstable_template.pptx")
        slide = ppt.slides[0]

        # Replace title
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and any(tok in (shape.text or "") for tok in ["{TBD Title}", "{TBD TITLE}"]):
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f"{base_title} (Page {i + 1} of {total_slides})"
                run.font.size = Pt(26)
                run.font.bold = True
                run.font.name = "Montserrat"
                from pptx.dml.color import RGBColor
                run.font.color.rgb = RGBColor(255, 255, 255)  # match header style
                p.alignment = PP_ALIGN.LEFT  # or CENTER if preferred
                break

        # Table anchor
        anchor = _find_anchor(slide, "tableanchor")
        if anchor is None:
            left, top, width = Inches(0.3), Inches(1.7), Inches(8.0)
        else:
            left, top, width = anchor.left, anchor.top, anchor.width

        rows = len(batch) + 1
        cols = 5
        table = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.8)).table
        table.columns[0].width = Inches(2.4)
        table.columns[1].width = Inches(2.4)
        table.columns[2].width = Inches(1.3)
        table.columns[3].width = Inches(1.0)
        table.columns[4].width = Inches(1.0)

        headers = ["Business Name", "Address", "Revenue", "YoY Growth", "Ticket Size"]
        # Header row
        for c, header in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = header
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(9)

        # Body rows (slightly smaller)
        for r, biz in enumerate(batch, start=1):
            cells = [
                biz.get("name", ""),
                biz.get("address", ""),
                f"${biz.get('annual_revenue', 0):,.0f}",
                f"{biz.get('yoy_growth', 0) * 100:+.1f}%",
                f"${biz.get('ticket_size', 0):,.0f}",
            ]
            for c, text in enumerate(cells):
                cell = table.cell(r, c)
                cell.text = text
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(8)
            # tighten row height a touch
            try:
                table.rows[r].height = Inches(0.25)
            except Exception:
                pass

        slide_number = 41 + i
        slide_path = os.path.join(output_dir, f"slide_{slide_number:02}_BusinessTable.pptx")
        ppt.save(slide_path)
        print(f"✅ Saved: {slide_path}")
