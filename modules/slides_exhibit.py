# slides_exhibit.py — template-driven exhibits
# Uses the downloaded exhibit PPTX template so all charts inherit the same
# header/footer/margins as the rest of the deck.

import matplotlib.pyplot as plt
from modules.map_generator import generate_map_png_from_summaries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import os

EXHIBIT_TEMPLATE = "modules/downloaded_exhibit_template.pptx"

# ------------------------
# Style for matplotlib PNGs
# ------------------------

import matplotlib.pyplot as plt

def apply_peerview_style():
    """Apply PeerView style; fallback if Montserrat isn't installed."""
    try:
        import matplotlib.font_manager as fm
        has_montserrat = any("Montserrat" in f.name for f in fm.fontManager.ttflist)
    except Exception:
        has_montserrat = False

    base_font = "Montserrat" if has_montserrat else "DejaVu Sans"

    plt.rcParams.update({
        "font.family": base_font,
        "axes.facecolor": "#FFFFFF",
        "figure.facecolor": "#FFFFFF",
        "axes.edgecolor": "#CCCCCC",
        "axes.titleweight": "bold",
        "axes.titlesize": 16,
        "axes.labelcolor": "#333333",
        "xtick.color": "#333333",
        "ytick.color": "#333333",
        "axes.grid": False,
    })


# ------------------------
# PPT helper utilities
# ------------------------

def _replace_in_runs(shape, mapping: dict) -> bool:
    """Replace substrings inside runs to preserve template formatting.
    Returns True if any replacement occurred.
    """
    if not getattr(shape, "has_text_frame", False):
        return False
    changed = False
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            if not r.text:
                continue
            new_text = r.text
            for k, v in mapping.items():
                if k in new_text:
                    new_text = new_text.replace(k, v)
                    changed = True
            r.text = new_text
    return changed


def _find_named(slide, *names):
    for shp in slide.shapes:
        if getattr(shp, "name", "") in names:
            return shp
    return None


def _chart_anchor(slide):
    """Find where to place the chart.
    Prefer a named anchor; else the largest rectangle; else fallback margins.
    Returns (left, top, width, height).
    """
    # 1) Named anchor(s)
    anchor = _find_named(slide, "ChartAnchor", "Chart", "ImageAnchor")
    if anchor:
        return anchor.left, anchor.top, anchor.width, anchor.height

    # 2) Largest rectangle
    max_area = 0
    best = None
    for shp in slide.shapes:
        try:
            if shp.shape_type == 1:  # rectangle
                area = shp.width * shp.height
                if area > max_area:
                    max_area = area
                    best = shp
        except Exception:
            pass
    if best:
        return best.left, best.top, best.width, best.height

    # 3) Fallback margins
    left = Inches(0.75)
    top = Inches(1.2)  # below the header bar
    width = slide.part.presentation.slide_width - 2 * left
    height = Inches(4.0)
    return left, top, width, height

# ------------------------
# Template-driven slide builder
# ------------------------

def generate_chart_slide(chart_title: str, image_path: str, summary_text: str) -> Presentation:
    """Open the exhibit template and fill: title, chart image, and analysis box.
    Keeps all typography and spacing from the template.
    """
    ppt = Presentation(EXHIBIT_TEMPLATE)
    slide = ppt.slides[0]

    # 1) Replace the header title placeholder (format-preserving)
    title_done = False
    for shape in slide.shapes:
        if _replace_in_runs(shape, {"{TBD EXHIBIT TITLE}": chart_title, "Exhibit {TBD}": chart_title}):
            title_done = True
            break
    if not title_done:
        # Fallback: first text frame gets the title
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                shape.text_frame.clear()
                shape.text_frame.paragraphs[0].text = chart_title
                break

    # 2) Place the chart image using the anchor logic
    left, top, width, height = _chart_anchor(slide)
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)

    # 3) Fill the analysis text into the box that contains {TBD ANALYSIS}
    analysis_shape = None
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False):
            if "{TBD ANALYSIS}" in (shp.text_frame.text or ""):
                analysis_shape = shp
                break

    if analysis_shape:
        tf = analysis_shape.text_frame
        # Preserve a heading "Analysis" if present, then inject body
        has_heading = any("Analysis" in (r.text or "") for p in tf.paragraphs for r in p.runs)
        tf.clear()
        if has_heading:
            tf.paragraphs[0].text = "Analysis"
        p = tf.add_paragraph()
        p.text = summary_text
        p.alignment = p.alignment or PP_ALIGN.LEFT
    else:
        # Last resort: bottom box within margins
        txtbox = slide.shapes.add_textbox(Inches(0.75), Inches(5.3), slide.part.presentation.slide_width - Inches(1.5), Inches(2.0))
        tf = txtbox.text_frame
        tf.paragraphs[0].text = "Analysis"
        p = tf.add_paragraph()
        p.text = summary_text

    return ppt

# ------------------------
# Chart generators (unchanged except for light style)
# ------------------------

def generate_revenue_chart(path, summaries, end_date: str):
    apply_peerview_style()

    # Sort trusted businesses by revenue
    trusted = [b for b in summaries if b.get("benchmark") == "trusted"]
    trusted = sorted(trusted, key=lambda x: x["annual_revenue"], reverse=True)

    # Disambiguate duplicate names
    seen_names = {}
    def disambiguate(name):
        base = name[:20]
        if base in seen_names:
            seen_names[base] += 1
            return f"{base[:17]}…{seen_names[base]}"
        seen_names[base] = 1
        return base if len(name) <= 20 else base[:19] + "…1"

    names = [disambiguate(b["name"]) for b in trusted]
    values = [b["annual_revenue"] for b in trusted]
    values_m = [v / 1_000_000 for v in values]

    mean_val = sum(values) / len(values)
    median_val = sorted(values)[len(values) // 2]

    colors = ["#D4AF37", "#C0C0C0", "#CD7F32"] + ["#A2D5AB"] * (len(values) - 3)

    fig, ax = plt.subplots(figsize=(12, 5.5))
    ax.bar(names, values_m, color=colors, width=0.6)

    for bar, val in zip(ax.patches, values_m):
        ax.text(bar.get_x() + bar.get_width()/2, val + 0.12, f"${val:.1f}M", ha="center", va="bottom", fontsize=8, rotation=90)

    ax.axhline(mean_val/1_000_000, color="#4682B4", linestyle="--", linewidth=1, label=f"Mean: ${mean_val/1_000_000:.1f}M")
    ax.axhline(median_val/1_000_000, color="#9370DB", linestyle=":", linewidth=1, label=f"Median: ${median_val/1_000_000:.1f}M")

    ax.set_title("Annual Revenue", fontsize=16, fontweight="bold", color="#333333", pad=20)
    if end_date:
        ax.text(0.5, 1.03, f"As of {end_date}", transform=ax.transAxes, fontsize=10, color="#555555", ha="center")

    ax.set_ylabel("Revenue ($M)", fontsize=11, color="#333333")
    ax.set_xticklabels(names, rotation=45, ha="right", fontsize=8, color="#333333")
    ax.tick_params(axis="y", labelsize=9, colors="#333333")
    ax.legend(loc="upper right", frameon=False)

    plt.tight_layout()
    plt.savefig(path)
    plt.close()
    return True


def generate_yoy_chart(path, summaries, end_date: str):
    apply_peerview_style()

    trusted = [b for b in summaries if b.get("benchmark") == "trusted" and b.get("yoy_growth") is not None]
    trusted = sorted(trusted, key=lambda x: x["yoy_growth"], reverse=True)

    seen_names = {}
    def disambiguate(name):
        base = name[:20]
        if base in seen_names:
            seen_names[base] += 1
            return f"{base[:17]}…{seen_names[base]}"
        seen_names[base] = 1
        return base if len(name) <= 20 else base[:19] + "…1"

    names = [disambiguate(b["name"]) for b in trusted]
    values = [round(b["yoy_growth"] * 100) for b in trusted]

    avg = sum(values) / len(values)
    median = sorted(values)[len(values) // 2]

    colors = ["#4CAF50" if v >= 0 else "#E57373" for v in values]

    fig, ax = plt.subplots(figsize=(12, 5.5))
    ax.bar(names, values, color=colors, width=0.6)

    for bar, val in zip(ax.patches, values):
        offset = 1.0 if abs(val) < 10 else 0.5
        ax.text(bar.get_x()+bar.get_width()/2, val + (offset if val >= 0 else -offset), f"{val:.0f}%",
                ha="center", va="bottom" if val >= 0 else "top", fontsize=8, weight="bold")

    ax.axhline(avg, color="#4682B4", linestyle="--", linewidth=1, label=f"Mean: {avg:.1f}%")
    ax.axhline(median, color="#9370DB", linestyle=":", linewidth=1, label=f"Median: {median:.1f}%")
    ax.axhline(0, color="#CCCCCC", linewidth=0.5)

    ax.set_title("Year over Year Revenue Growth", fontsize=16, fontweight="bold", color="#333333", pad=28)
    if end_date:
        ax.text(0.5, 1.06, f"As of {end_date}", transform=ax.transAxes, fontsize=10, color="#555555", ha="center")

    ax.set_ylabel("Growth (%)", fontsize=11, color="#333333")
    ax.set_xticklabels(names, rotation=45, ha="right", fontsize=8, color="#333333")
    ax.tick_params(axis="y", labelsize=9, colors="#333333")
    ax.legend(loc="upper right", frameon=False)

    plt.tight_layout()
    plt.savefig(path)
    plt.close()
    return True


def generate_ticket_chart(path, summaries, end_date: str):
    apply_peerview_style()

    trusted = [b for b in summaries if b.get("benchmark") == "trusted" and b.get("ticket_size") is not None]
    trusted = sorted(trusted, key=lambda x: x["ticket_size"], reverse=True)

    seen_names = {}
    def disambiguate(name):
        base = name[:20]
        if base in seen_names:
            seen_names[base] += 1
            return f"{base[:17]}…{seen_names[base]}"
        seen_names[base] = 1
        return base if len(name) <= 20 else base[:19] + "…1"

    names = [disambiguate(b["name"]) for b in trusted]
    values = [round(b["ticket_size"]) for b in trusted]

    mean_val = sum(values) / len(values)
    median_val = sorted(values)[len(values) // 2]

    fig, ax = plt.subplots(figsize=(12, 5.5))
    ax.bar(names, values, color="#4CAF50", width=0.6)

    for bar, val in zip(ax.patches, values):
        ax.text(bar.get_x()+bar.get_width()/2, val + 0.5, f"${val}", ha="center", va="bottom", fontsize=8, weight="bold")

    ax.axhline(mean_val, color="#4682B4", linestyle="--", linewidth=1, label=f"Mean: ${mean_val:.0f}")
    ax.axhline(median_val, color="#9370DB", linestyle=":", linewidth=1, label=f"Median: ${median_val:.0f}")
    ax.axhline(0, color="#CCCCCC", linewidth=0.5)

    ax.set_title("Average Ticket Size", fontsize=16, fontweight="bold", color="#333333", pad=28)
    if end_date:
        ax.text(0.5, 1.06, f"As of {end_date}", transform=ax.transAxes, fontsize=10, color="#555555", ha="center")

    ax.set_ylabel("Dollars ($)", fontsize=11, color="#333333")
    ax.set_xticklabels(names, rotation=45, ha="right", fontsize=8, color="#333333")
    ax.tick_params(axis="y", labelsize=9, colors="#333333")
    ax.legend(loc="upper right", frameon=False)

    plt.tight_layout()
    plt.savefig(path)
    plt.close()
    return True


def generate_market_size_chart(path, summaries, end_date: str):
    """
    Bar chart with dynamic headroom so value labels never collide with bars/title.
    Self-contained: uses project style if available, else a safe default.
    """
    # Try to apply the shared style; fall back to a minimal one if not present
    try:
        apply_peerview_style()  # defined in slides_exhibit.py in this project
    except Exception:
        plt.rcParams.update({
            "font.family": "Montserrat",
            "axes.facecolor": "#FFFFFF",
            "figure.facecolor": "#F8F8F8",
            "axes.edgecolor": "#CCCCCC",
            "axes.grid": False,
            "xtick.color": "#333333",
            "ytick.color": "#333333",
        })

    # Data split
    trusted = [b for b in summaries if b.get("benchmark") == "trusted" and b.get("annual_revenue") is not None]
    # If your untrusted flag is something else, swap "low" as needed:
    untrusted = [b for b in summaries if b.get("benchmark") == "low" and b.get("annual_revenue") is not None]

    trusted_total = sum(b["annual_revenue"] for b in trusted)
    num_trusted = len(trusted)
    num_untrusted = len(untrusted)
    projected_total = trusted_total * (num_trusted + num_untrusted) / max(num_trusted, 1)

    lower_m = trusted_total / 1_000_000
    upper_m = projected_total / 1_000_000

    # Plot
    fig, ax = plt.subplots(figsize=(8, 5))
    bars = ax.bar(["Verified Revenue", "Projected Total"], [lower_m, upper_m],
                  color=["#4CAF50", "#C0C0C0"], edgecolor="black", width=0.5)
    bars[1].set_hatch("//")  # distinguish projected even in grayscale prints

    # Subtle frame
    ax.set_facecolor("#FFFFFF")
    for spine in ax.spines.values():
        spine.set_visible(True)
        spine.set_linewidth(1)
        spine.set_edgecolor("#CCCCCC")

    # --- Fix label overlap: add headroom and use dynamic offsets ---
    vals = [lower_m, upper_m]
    ymax = max(vals) if max(vals) > 0 else 1
    ax.set_ylim(0, ymax * 1.25)     # ~25% headroom above tallest bar
    yoffset = 0.04 * ax.get_ylim()[1]  # offset scales with axis height

    for bar, val in zip(bars, vals):
        x = bar.get_x() + bar.get_width() / 2
        y = bar.get_height() + yoffset
        ax.annotate(f"${val:.1f}M", (x, y),
                    ha="center", va="bottom", fontsize=9, weight="bold", clip_on=False)

    # Business counts — small offset so these don’t collide either
    ax.text(0, lower_m + yoffset/2, f"{num_trusted} businesses", ha='center', fontsize=8, color="#333333")
    ax.text(1, upper_m + yoffset/2, f"{num_trusted + num_untrusted} total (incl. {num_untrusted} projected)",
            ha='center', fontsize=8, color="#333333")

    # Axes/title
    ax.axhline(0, color="#CCCCCC", linewidth=0.5)
    ax.set_ylabel("Revenue ($M)", fontsize=11, color="#333333")
    ax.set_title("Estimated Market Revenue Potential", fontsize=16, fontweight='bold', color="#333333", pad=28)
    if end_date:
        ax.text(0.5, 1.06, f"As of {end_date}", transform=ax.transAxes, fontsize=10, color="#555555", ha='center')

    ax.set_xticks([0, 1])
    ax.set_xticklabels(["Verified Revenue", "Projected Total"], fontsize=9, color="#333333")
    ax.tick_params(axis='y', labelsize=9, colors="#333333")
    ax.margins(y=0.10)

    plt.tight_layout()
    plt.savefig(path)
    plt.close()
    return True


def get_market_size_analysis():
    return (
        "This chart compares the total verified revenue from businesses with high-quality data to an estimate "
        "for the full local market. The upper bound assumes that businesses without usable data perform similarly "
        "to those with data, which likely overstates true market size. Poor data quality is often associated with smaller "
        "businesses or those facing operational challenges."
    )


from pptx import Presentation
from pptx.util import Inches

# Helper copied from your file (lightweight):


def _find_named(slide, *names):
    for shp in slide.shapes:
        if getattr(shp, "name", "") in names:
            return shp
    return None

def _chart_anchor_dims_from_template(template_path: str):
    ppt = Presentation(template_path)
    slide = ppt.slides[0]
    anchor = _find_named(slide, "ChartAnchor", "Chart", "ImageAnchor")
    if anchor:
        return anchor.left, anchor.top, anchor.width, anchor.height
    # fallback = largest rectangle
    max_area, best = 0, None
    for shp in slide.shapes:
        try:
            area = shp.width * shp.height
            if area > max_area:
                max_area, best = area, shp
        except Exception:
            pass
    return (best.left, best.top, best.width, best.height) if best else (None, None, None, None)

def generate_map_chart(output_path, summaries):
    try:
        from modules.slides_exhibit import EXHIBIT_TEMPLATE
    except Exception:
        EXHIBIT_TEMPLATE = "modules/downloaded_exhibit_template.pptx"
    _, _, width, height = _chart_anchor_dims_from_template(EXHIBIT_TEMPLATE)
    aspect_ratio = float(width) / float(height) if height else (3/2)
    from modules.map_generator import generate_map_png_from_summaries
    return generate_map_png_from_summaries(
        summaries,
        output_path,
        zoom_fraction=0.75,
        aspect_ratio=aspect_ratio,
        window_height_px=800,
    )


def build_exhibit_slide_from_template(chart_png_path: str, exhibit_title: str, analysis_text: str,
                                      template_path: str = EXHIBIT_TEMPLATE) -> Presentation:
    """Preferred entrypoint if you want direct control.
    Kept for backwards compatibility with earlier code.
    """
    return generate_chart_slide(exhibit_title, chart_png_path, analysis_text)
